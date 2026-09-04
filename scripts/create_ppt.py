from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import os

def create_pitch_deck():
    prs = Presentation()
    
    # 1. Title Slide
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "NETRA"
    subtitle.text = "Neural Evaluation & Tracking Research Architecture\nMulti-Modal Deepfake & Scam Detection"

    # 2. The Problem
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    body = slide.placeholders[1]
    title.text = "The Problem"
    tf = body.text_frame
    tf.text = "Deepfakes and scams are destroying trust."
    p = tf.add_paragraph()
    p.text = "- Existing tools are 'black boxes' with no explanation."
    p = tf.add_paragraph()
    p.text = "- Financial scams using AI voice clones are skyrocketing."
    p = tf.add_paragraph()
    p.text = "- India lacks localized, explainable forensic tools."

    # 3. The Solution
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    body = slide.placeholders[1]
    title.text = "The NETRA Solution"
    tf = body.text_frame
    tf.text = "A transparent, multi-modal forensic platform."
    p = tf.add_paragraph()
    p.text = "1. Video Face-Swap Detection (Spatial/Frequency Analysis)"
    p = tf.add_paragraph()
    p.text = "2. Audio Voice Clone Detection"
    p = tf.add_paragraph()
    p.text = "3. Financial Scam & Phishing Detection (Text/URL)"

    # 4. The Architecture (AWS + Kaggle)
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    body = slide.placeholders[1]
    title.text = "Hybrid Cloud Architecture"
    tf = body.text_frame
    tf.text = "Optimized for Cost, Speed, and Scalability."
    p = tf.add_paragraph()
    p.text = "• Kaggle GPUs (P100/T4): Used for 100% of Heavy ML Training."
    p = tf.add_paragraph()
    p.text = "• AWS EC2/Lambda: Used for lightweight, fast inference."
    p = tf.add_paragraph()
    p.text = "• Forensic Report Engine: 100% Offline deterministic legal synthesis."

    # 5. How it Works (The Forensic Engine)
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    body = slide.placeholders[1]
    title.text = "The Forensic Dossier Engine"
    tf = body.text_frame
    tf.text = "NETRA doesn't just say 'Fake'. It explains WHY."
    p = tf.add_paragraph()
    p.text = "1. ML Detectors extract raw evidence (JSON)."
    p = tf.add_paragraph()
    p.text = "2. Deterministic Engine v5.0 compiles verified telemetry."
    p = tf.add_paragraph()
    p.text = "3. Generates court-admissible Section 65B forensic reports."

    # 6. The Scam Detector Module
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    body = slide.placeholders[1]
    title.text = "The Scam Detector"
    tf = body.text_frame
    tf.text = "Catching fraud before it happens."
    p = tf.add_paragraph()
    p.text = "- High-throughput approach: 100+ Hardcoded Rules + Random Forest Classifier."
    p = tf.add_paragraph()
    p.text = "- Identifies phishing links, urgency markers, and crypto scams."

    # 7. The User Experience
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    body = slide.placeholders[1]
    title.text = "The User Experience"
    tf = body.text_frame
    tf.text = "Designed for Investigators and Journalists."
    p = tf.add_paragraph()
    p.text = "• Interactive Evidence Timeline: Scrub through video and see red flags."
    p = tf.add_paragraph()
    p.text = "• Detailed probability graphs for spatial/frequency anomalies."

    # 8. Roadmap & Conclusion
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    body = slide.placeholders[1]
    title.text = "Roadmap & Future Impact"
    tf = body.text_frame
    tf.text = "NETRA is built for scale."
    p = tf.add_paragraph()
    p.text = "- Zero-budget baseline architecture."
    p = tf.add_paragraph()
    p.text = "- Open-source API for fact-checking organizations."
    
    prs.save('/Users/iamsparsh00321/Desktop/NETRA_pitch_deck_v6.pptx')
    print("Successfully created NETRA_pitch_deck_v6.pptx")

if __name__ == '__main__':
    create_pitch_deck()
